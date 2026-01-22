from __future__ import annotations

from pathlib import Path
from typing import List

from pptx import Presentation

from ..models import SlideChunk


def _collect_text(shape) -> str:
    if hasattr(shape, "text"):
        return shape.text
    if shape.has_table:
        rows = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)
    return ""


def parse_ppt(ppt_path: Path) -> List[SlideChunk]:
    prs = Presentation(str(ppt_path))
    slides: List[SlideChunk] = []
    current_section: str | None = None
    for idx, slide in enumerate(prs.slides, start=1):
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        title = title_shapes[0].text if title_shapes else None
        bullets: List[str] = []
        levels: List[int] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            paragraphs = []
            para_levels = []
            for paragraph in shape.text_frame.paragraphs:
                fragments = [run.text for run in paragraph.runs if run.text]
                text_line = "".join(fragments).strip()
                if text_line:
                    paragraphs.append(text_line)
                    para_levels.append(paragraph.level)
            if paragraphs:
                bullets.extend(paragraphs)
                levels.extend(para_levels)
        raw_text = "\n".join(bullets)
        notes = slide.has_notes_slide and slide.notes_slide.notes_text_frame.text or None
        is_heading = bool(title and not raw_text.strip())
        if is_heading:
            current_section = title
        slides.append(
            SlideChunk(
                slide_number=idx,
                title=title,
                bullets=bullets,
                notes=notes,
                raw_text=raw_text,
                section=current_section,
                is_heading=is_heading,
                level=min(levels) if levels else 0,
            )
        )
    last_section = None
    for slide in slides:
        if slide.is_heading:
            last_section = slide.title
        if not slide.section:
            slide.section = last_section
    return slides
